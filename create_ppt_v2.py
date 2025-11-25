#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors matching HTML
NAVY = RGBColor(15, 23, 42)
GREEN = RGBColor(16, 185, 129)
BLUE = RGBColor(59, 130, 246)
RED = RGBColor(239, 68, 68)
ORANGE = RGBColor(245, 158, 11)
GRAY = RGBColor(51, 65, 85)
LIGHT_GRAY = RGBColor(148, 163, 184)

def add_title_slide(prs, subtitle_small, title, subtitle, footer_text=""):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Navy background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Small subtitle at top
    if subtitle_small:
        small_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.4))
        small_frame = small_box.text_frame
        small_frame.text = subtitle_small
        p = small_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.LEFT

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(70)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.alignment = PP_ALIGN.LEFT

    # Footer
    if footer_text:
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(6), Inches(1))
        footer_frame = footer_box.text_frame
        for i, line in enumerate(footer_text.split('\n')):
            if i > 0:
                footer_frame.add_paragraph()
            p = footer_frame.paragraphs[i]
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = LIGHT_GRAY
            p.alignment = PP_ALIGN.LEFT

def add_content_slide(prs, title):
    """일반 콘텐츠 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title with underline
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Green underline
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1), Inches(9.5), Inches(1))
    line.line.color.rgb = GREEN
    line.line.width = Pt(4)

    return slide

# Slide 1: 표지
add_title_slide(
    prs,
    "비영리조직 사업(창업)계획서",
    "Re:Spring",
    "취약계층의 자립을 돕는 친환경 업사이클링 제조 솔루션",
    "과목: 비영리조직 창업 및 운영\n제출자: 사회복지학과 202531606 민소은"
)

# Slide 2: 회사소개
slide = add_content_slide(prs, "① 회사소개: 기관의 비전, 목적")
vision = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(7), Inches(0.6))
vision.text_frame.text = '"환경과 복지의 교차점에서 지속 가능한 일자리를 창출한다"'
vision.text_frame.paragraphs[0].font.size = Pt(22)
vision.text_frame.paragraphs[0].font.bold = True
vision.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

y = 2.8
items = [
    ("설립 목적", "단순 현금 지원이 아닌 '제조업 일자리'를 통해 취약계층의 경제적/정서적 완전 자립 실현."),
    ("기관 비전 (2030)", "지역사회 1호 '장애인 표준사업장 인증' 친환경 제조 전문 기업 도약."),
    ("핵심 가치", "치유(Healing)\n전문(Pro)\n순환(Eco)")
]
for i, (h, t) in enumerate(items):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.7), Inches(3.2))
    tf = box.text_frame
    tf.text = h
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = t
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    for p in tf.paragraphs[2:]:
        p.font.size = Pt(13)

# Slide 3: 사업소개
slide = add_content_slide(prs, "② 사업소개: 창업배경 및 창업동기")
box1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.2), Inches(4))
tf1 = box1.text_frame
tf1.text = "창업 배경 (문제인식)"
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "• 복지 사각지대: 관내 발달장애인/노인 취업률 20% 미만. 기존 단순 임가공 일자리는 '지속 가능한 급여' 지급 불가."
tf1.add_paragraph().text = "• 환경 위기: 플라스틱 소각 비용 급증 및 기업 ESG 실적 압박 심화."
tf1.paragraphs[0].font.size = Pt(17)
tf1.paragraphs[0].font.bold = True
for p in tf1.paragraphs[2:]:
    p.font.size = Pt(13)

box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.2), Inches(4))
tf2 = box2.text_frame
tf2.text = "창업 동기"
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = '"45세 사회복지 전문가로서 현장의 한계를 절감했습니다."'
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "후원금에 의존하는 복지는 지속 가능하지 않습니다. '제품 경쟁력'으로 당당하게 월급을 주는 기업을 만들기 위해 창업을 결심했습니다."
tf2.paragraphs[0].font.size = Pt(17)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[2].font.size = Pt(15)
tf2.paragraphs[2].font.bold = True
for p in tf2.paragraphs[4:]:
    p.font.size = Pt(13)

# Slide 4: 비즈니스 모델
slide = add_content_slide(prs, "② 비즈니스 모델 (온/오프라인, 글로벌)")
subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8), Inches(0.3))
subtitle.text_frame.text = "안정적인 수익 구조 확보를 위한 다각화 전략"
subtitle.text_frame.paragraphs[0].font.size = Pt(14)

models = [
    ("1. 오프라인 (B2G/B2B)", "대상: 관공서, 보건소, 대기업 ESG팀\n• '중증장애인생산품 우선구매' 활용 수의계약\n• 기업 사내 캠페인 + 굿즈 납품"),
    ("2. 온라인 (D2C)", "대상: 가치소비 MZ세대\n• 자사몰 및 스마트스토어 운영\n• 와디즈/텀블벅 펀딩을 통한 신제품 런칭"),
    ("3. 글로벌 사업 (Global)", "대상: 베트남 등 동남아 시장\n• 5년 차 해외 판로 개척\n• 한국형 업사이클링 디자인(K-Eco) 수출")
]
for i, (h, t) in enumerate(models):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.7), Inches(3.5))
    tf = box.text_frame
    tf.text = h
    tf.add_paragraph().text = ""
    for line in t.split('\n'):
        tf.add_paragraph().text = line
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    for p in tf.paragraphs[2:]:
        p.font.size = Pt(12)

# Slide 5: 제품 소개
slide = add_content_slide(prs, "핵심 제품 소개 (Signature Products)")
box1 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(4), Inches(3.5))
tf1 = box1.text_frame
tf1.text = "1. 마블링 화분"
tf1.add_paragraph().text = "폐플라스틱 고유의 패턴을 살린 세상에 하나뿐인 디자인.\n기업 로고 각인 서비스 제공."
tf1.add_paragraph().text = "예상가: 12,000원"
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "2. 아웃도어 카라비너"
tf1.add_paragraph().text = "고강도 플라스틱(HDPE)을 활용한 캠핑/등산용 굿즈.\n내구성 테스트 완료."
tf1.add_paragraph().text = "예상가: 3,500원"
tf1.paragraphs[0].font.size = Pt(15)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[2].font.color.rgb = GREEN
tf1.paragraphs[2].font.bold = True
tf1.paragraphs[4].font.size = Pt(15)
tf1.paragraphs[4].font.bold = True
tf1.paragraphs[6].font.color.rgb = BLUE
tf1.paragraphs[6].font.bold = True
for p in [tf1.paragraphs[1], tf1.paragraphs[5]]:
    p.font.size = Pt(13)

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2.2), Inches(3.8), Inches(3))
tf2 = box2.text_frame
tf2.text = "제품 차별성"
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "자체 내구성 테스트 완료, 고품질 원료 사용으로 저가 제품 대비 우수한 품질."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "제품 하단 'Made by 000(근로자 실명)' 각인 서비스."
tf2.paragraphs[0].font.size = Pt(17)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[2].font.size = Pt(13)
tf2.paragraphs[4].font.size = Pt(13)
tf2.paragraphs[4].font.bold = True

# Slide 6: SWOT
slide = add_content_slide(prs, "④ 시장분석: SWOT 분석")
swot = [
    ("Strength (강점)", "• 복지 전문가(대표) + 기술 전문가(공장장) 시너지\n• 정부 인건비 지원 수혜 (가격 경쟁력 확보)\n• 독창적인 금형 디자인 보유", BLUE),
    ("Weakness (약점)", "• 초기 브랜드 인지도 부족\n• 전용 설비(사출기) 구축 초기 투자비용 부담", RED),
    ("Opportunity (기회)", "• 공공기관 우선구매 시장(20조원) 확대\n• 기업 ESG 경영 강화로 친환경 굿즈 수요 폭증\n• 가치소비 트렌드 확산", GREEN),
    ("Threat (위협)", "• 저가 중국산 제품과의 가격 경쟁\n• 유사 업사이클링 업체의 난립", ORANGE)
]
pos = [(0.8, 1.8), (5.2, 1.8), (0.8, 4.2), (5.2, 4.2)]
for (h, t, c), (x, y) in zip(swot, pos):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(2))
    tf = box.text_frame
    tf.text = h
    tf.add_paragraph().text = ""
    for line in t.split('\n'):
        tf.add_paragraph().text = line
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = c
    for p in tf.paragraphs[2:]:
        p.font.size = Pt(12)

# Slide 7: 마케팅
slide = add_content_slide(prs, "④ 마케팅 전략 및 판로 제시")
box1 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(4), Inches(3.5))
tf1 = box1.text_frame
tf1.text = "홍보 전략"
tf1.add_paragraph().text = '"제품이 아닌 가치를 팝니다"'
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "• 스토리텔링: 제품 QR코드로 제작자(장애인)의 작업 영상 연결."
tf1.add_paragraph().text = "• 체험단 운영: 지역 맘카페 연계 '나만의 화분 만들기' 체험단."
tf1.paragraphs[0].font.size = Pt(16)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[1].font.size = Pt(14)
tf1.paragraphs[1].font.bold = True
for p in tf1.paragraphs[3:]:
    p.font.size = Pt(12)

box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(4.5))
tf2 = box2.text_frame
tf2.text = "단계별 판로 개척"
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "1단계 (공공): 나라장터, 꿈드래 쇼핑몰 입점."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "2단계 (온라인): 네이버 스마트스토어, 와디즈 펀딩."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "3단계 (제휴): 제로웨이스트 샵 20개소 입점."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "4단계 (수출): KOTRA 연계 해외 전시회 참가."
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.bold = True
for i in [2, 4, 6, 8]:
    tf2.paragraphs[i].font.size = Pt(12)

# Slide 8: 인력구성
slide = add_content_slide(prs, "③ 사업계획: 인력구성 (HR)")
roles = [
    ("대표 민소은 (본인)", "사회복지 전문가 (1급)", "인력 관리 및 직무 지도, 지자체/관공서 영업 총괄."),
    ("기술 이사 (CTO)", "생산 총괄 (경력 20년)", "사출 금형 설계, 공장장 출신, 생산 라인 및 품질 관리."),
    ("현장 근로자 (3명)", "취약계층 우선 채용", "경력단절여성 및 장애인, 지역 사회복지관 추천 채용.")
]
for i, (name, role, desc) in enumerate(roles):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.7), Inches(3.5))
    tf = box.text_frame
    tf.text = name
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = role
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = desc
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[2].font.size = Pt(13)
    tf.paragraphs[2].font.bold = True
    tf.paragraphs[4].font.size = Pt(11)

# Slide 9: 임금 계획
slide = add_content_slide(prs, "③ 사업계획: 인력 운영 및 임금 (2025년 기준)")
table_data = [
    ["구분", "대상", "월 급여 (세전)", "비고"],
    ["대표", "대표 민소은", "무급", "초기 3년간 무급 (재투자)"],
    ["관리직", "기술이사", "3,000,000원", "경력직 대우"],
    ["현장직", "취약계층 근로자", "2,096,270원", "최저임금 준수"]
]
table = slide.shapes.add_table(4, 4, Inches(1), Inches(2), Inches(8), Inches(2)).table
for i, row in enumerate(table_data):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            if i == 0:
                p.font.bold = True

note = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1.2))
tf = note.text_frame
tf.text = "※ 임금 산출 근거: 2025년 최저시급 10,030원 × 209시간 (주휴수당 포함 월 소정근로시간)"
tf.add_paragraph().text = "※ 복리후생: 4대보험 가입(두루누리 활용), 식대 별도 제공, 심리상담(EAP) 프로그램 지원"
for p in tf.paragraphs:
    p.font.size = Pt(11)

# Slide 10: 정부 제도
slide = add_content_slide(prs, "③ 사업계획: 정부 제도 활용 전략 (필수 기입)")
stages = [
    ("1. 진입기", "사회적기업가 육성사업", "창업 초기 자금 3,000만원 확보하여 핵심 설비(사출기) 구입 비용으로 사용."),
    ("2. 성장기", "일자리창출사업 (예비)", "예비사회적기업 지정 후, 취약계층 신규 채용 인건비의 50~70%를 지원받아 고정비 절감."),
    ("3. 도약기", "장애인 표준사업장 지원", "한국장애인고용공단 시설 무상지원금을 활용하여 작업장 환경 개선.")
]
for i, (stage, title, desc) in enumerate(stages):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.7), Inches(3.5))
    tf = box.text_frame
    tf.text = stage
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = title
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = desc
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[2].font.size = Pt(14)
    tf.paragraphs[2].font.bold = True
    tf.paragraphs[4].font.size = Pt(11)

# Slide 11: 초기 투자
slide = add_content_slide(prs, "초기 투자 및 생산 설비 계획")
box1 = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(3.5), Inches(3.5))
tf1 = box1.text_frame
tf1.text = "총 소요 예산: 1억 원"
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "• 시설자금 (6,000만원):\n  사출기, 분쇄기, 금형 제작비, 공장 보증금."
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "• 운전자금 (4,000만원):\n  초기 6개월 인건비, 재료비, 시제품 홍보비."
tf1.paragraphs[0].font.size = Pt(16)
tf1.paragraphs[0].font.bold = True
for p in tf1.paragraphs[2:]:
    p.font.size = Pt(13)

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3.5), Inches(3.5))
tf2 = box2.text_frame
tf2.text = "생산 설비 특징"
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 소형 사출기:\n  다품종 소량 생산 용이, 장애인 접근성 고려."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 저소음 분쇄기:\n  작업자 청력 보호를 위한 방음 박스 및 안전 센서 부착."
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.bold = True
for p in tf2.paragraphs[2:]:
    p.font.size = Pt(13)

# Slide 12: 자금 조달
slide = add_content_slide(prs, "자금 조달 방안")
subtitle = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(0.4))
subtitle.text_frame.text = "전략: 상환 부담 없는 정부지원금 비율(50%) 확대"
subtitle.text_frame.paragraphs[0].font.size = Pt(15)
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

funds = [
    ("자기자본 (20%)", "2,000만 원", "대표자 출자 (책임경영)"),
    ("정부지원 (50%)", "5,000만 원", "사회적기업 육성사업\n장애인공단 기기지원"),
    ("정책융자 (30%)", "3,000만 원", "서민금융진흥원 (저리)")
]
for i, (title, amt, desc) in enumerate(funds):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2.8), Inches(2.7), Inches(3))
    tf = box.text_frame
    tf.text = title
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = amt
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = desc
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[2].font.size = Pt(20)
    tf.paragraphs[2].font.bold = True
    tf.paragraphs[2].font.color.rgb = GREEN
    tf.paragraphs[4].font.size = Pt(12)

# Slide 13: 5년 계획
slide = add_content_slide(prs, "⑤ 향후 추진 계획 (5년 계획)")
years = [
    ("1년차 (진입)", "• 법인 설립\n• 예비사회적기업 지정\n• 시제품 5종 개발", BLUE),
    ("2년차 (성장)", "• 매출 3억 달성\n• 장애인 표준사업장 신청 준비\n• 고정 거래처 10곳", BLUE),
    ("3년차 (도약)", "• 사회적기업 본인증\n• 장애인 표준사업장 인증 완료\n• 공장 확장 이전\n• 취약계층 10명 고용", GREEN),
    ("4년차 (확장)", "• 신제품 라인업 확대\n• 지역 협력 공장 네트워크 구축\n• 전국 유통망 확대", GREEN),
    ("5년차 (Global)", "• 해외 수출 시작(베트남)\n• 연 매출 10억 달성\n• K-Eco 브랜드화", ORANGE)
]
for i, (title, content, color) in enumerate(years):
    x = 0.5 + (i * 1.9)
    box = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(1.7), Inches(4.5))
    tf = box.text_frame
    tf.text = title
    tf.add_paragraph().text = ""
    for line in content.split('\n'):
        tf.add_paragraph().text = line
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    for p in tf.paragraphs[2:]:
        p.font.size = Pt(10)

# Slide 14: 손익계산서
slide = add_content_slide(prs, "추정 손익 계산서 (3개년)")
financial = [
    ["구분 (단위: 백만 원)", "1년차", "2년차", "3년차"],
    ["매출액", "120", "300", "500"],
    ["매출원가", "40", "100", "165"],
    ["판관비 (인건비 포함)", "100", "180", "250"],
    ["영업이익", "△20", "20", "85"]
]
table = slide.shapes.add_table(5, 4, Inches(1.5), Inches(2.2), Inches(7), Inches(2.5)).table
for i, row in enumerate(financial):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            if i == 0 or i == 4:
                p.font.bold = True

note = slide.shapes.add_textbox(Inches(1.5), Inches(5.3), Inches(7), Inches(0.4))
note.text_frame.text = "* 1차년도 적자는 정부지원금(영업외수익)으로 보전하여 현금 흐름 유지."
note.text_frame.paragraphs[0].font.size = Pt(11)

# Slide 15: BEP
slide = add_content_slide(prs, "⑤ 향후 추진 계획: 손익분기점 (BEP)")
box1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(3.5), Inches(2.5))
tf1 = box1.text_frame
tf1.text = "BEP 달성 예상 시점"
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "8개월"
tf1.paragraphs[0].font.size = Pt(16)
tf1.paragraphs[2].font.size = Pt(60)
tf1.paragraphs[2].font.bold = True
tf1.paragraphs[2].font.color.rgb = GREEN
tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
tf1.paragraphs[2].alignment = PP_ALIGN.CENTER

box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(4.5))
tf2 = box2.text_frame
tf2.text = "달성 조건 및 전략"
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 달성 조건: 월 매출 1,400만 원\n  (BEP = 고정비 700만원 ÷ 공헌이익률 50%)."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 고정비 최소화: 월 700만 원\n  (인건비 지원 반영 후 자부담분)."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 공헌이익률: 50% (제조업 특성)."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 전략: B2B 대형 계약(월 500만 원 x 2건)\n  수주 시 조기 달성 가능."
tf2.paragraphs[0].font.size = Pt(15)
tf2.paragraphs[0].font.bold = True
for p in tf2.paragraphs[2:]:
    p.font.size = Pt(12)

# Slide 16: SROI
slide = add_content_slide(prs, "사회적 가치 측정 (SROI)")
box1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(3.5), Inches(3.5))
tf1 = box1.text_frame
tf1.text = "정량적 성과 (Quantitative)"
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "• 연간 폐플라스틱 25톤 재활용\n  (3년차 기준)."
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "• 탄소 감축 효과\n  (소나무 1,200그루 식재 상당)."
tf1.add_paragraph().text = ""
tf1.add_paragraph().text = "• 취약계층 10명 일자리 창출."
tf1.paragraphs[0].font.size = Pt(15)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[0].font.color.rgb = GREEN
for p in tf1.paragraphs[2:]:
    p.font.size = Pt(13)

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2.2), Inches(3.5), Inches(3.5))
tf2 = box2.text_frame
tf2.text = "정성적 성과 (Qualitative)"
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 근로자 우울감 척도 30% 개선."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 경제적 자립을 통한 탈수급\n  (수급자 탈피)."
tf2.add_paragraph().text = ""
tf2.add_paragraph().text = "• 지역사회 '환경 인식' 개선."
tf2.paragraphs[0].font.size = Pt(15)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = BLUE
for p in tf2.paragraphs[2:]:
    p.font.size = Pt(13)

# Slide 17: Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = NAVY

title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
title_box.text_frame.text = "Closing"
p = title_box.text_frame.paragraphs[0]
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.LEFT

msg_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3.5))
tf = msg_box.text_frame
tf.text = '"Re:Spring은 단순한 공장이 아닙니다.\n우리 지역사회의 가장 아픈 곳을 치유하는\n\'생산적 복지\'의 거점입니다."'
tf.add_paragraph().text = ""
tf.add_paragraph().text = "준비된 대표, 검증된 기술, 확실한 수요처.\n이제 지원금이라는 마중물이 필요합니다."
for p in tf.paragraphs:
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(226, 232, 240)
    p.alignment = PP_ALIGN.LEFT

# Save
prs.save('/home/user/claude-quickstarts/ReSpring_Business_Plan.pptx')
print("✓ PPT 완성! (HTML 내용 그대로 반영)")
