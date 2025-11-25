#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
NAVY = RGBColor(15, 23, 42)
GREEN = RGBColor(16, 185, 129)
BLUE = RGBColor(59, 130, 246)
RED = RGBColor(239, 68, 68)
ORANGE = RGBColor(245, 158, 11)
GRAY = RGBColor(51, 65, 85)

def add_title_slide(prs, title, subtitle, footer_text=""):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.alignment = PP_ALIGN.LEFT

    # Footer
    if footer_text:
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
        footer_frame = footer_box.text_frame
        footer_frame.text = footer_text
        p = footer_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(148, 163, 184)
        p.alignment = PP_ALIGN.LEFT

def add_content_slide(prs, title, content_dict):
    """일반 콘텐츠 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Title underline
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(9), Inches(0))
    line.line.color.rgb = GREEN
    line.line.width = Pt(4)

    return slide

# Slide 1: 표지
add_title_slide(
    prs,
    "Re:Spring",
    "취약계층의 자립을 돕는 친환경 업사이클링 제조 솔루션",
    "과목: 비영리조직 창업 및 운영\n제출자: 사회복지학과 202531606 민소은"
)

# Slide 2: 회사소개
slide = add_content_slide(prs, "① 회사소개: 기관의 비전, 목적", {})
vision_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
vision_frame = vision_box.text_frame
vision_frame.text = '"환경과 복지의 교차점에서 지속 가능한 일자리를 창출한다"'
p = vision_frame.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

y_pos = 2.8
texts = [
    ("설립 목적", "단순 현금 지원이 아닌 '제조업 일자리'를 통해 취약계층의 경제적/정서적 완전 자립 실현."),
    ("기관 비전 (2030)", "지역사회 1호 '장애인 표준사업장 인증' 친환경 제조 전문 기업 도약."),
    ("핵심 가치", "치유(Healing)\n전문(Pro)\n순환(Eco)")
]

for i, (label, text) in enumerate(texts):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(y_pos), Inches(2.8), Inches(2.5))
    frame = box.text_frame
    frame.text = f"{label}\n\n{text}"
    frame.paragraphs[0].font.size = Pt(16)
    frame.paragraphs[0].font.bold = True
    frame.paragraphs[0].font.color.rgb = NAVY
    for para in frame.paragraphs[1:]:
        para.font.size = Pt(14)

# Slide 3: 사업소개
slide = add_content_slide(prs, "② 사업소개: 창업배경 및 창업동기", {})

box1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4), Inches(2.5))
frame1 = box1.text_frame
frame1.text = """창업 배경 (문제인식)

• 복지 사각지대: 관내 발달장애인/노인 취업률 20% 미만
• 환경 위기: 플라스틱 소각 비용 급증"""
frame1.paragraphs[0].font.size = Pt(18)
frame1.paragraphs[0].font.bold = True
for para in frame1.paragraphs[1:]:
    para.font.size = Pt(14)

box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(2.5))
frame2 = box2.text_frame
frame2.text = """창업 동기

"45세 사회복지 전문가로서 현장의 한계를 절감했습니다."

후원금에 의존하는 복지는 지속 가능하지 않습니다. '제품 경쟁력'으로 당당하게 월급을 주는 기업을 만들기 위해 창업을 결심했습니다."""
frame2.paragraphs[0].font.size = Pt(18)
frame2.paragraphs[0].font.bold = True
for para in frame2.paragraphs[1:]:
    para.font.size = Pt(14)

# Slide 4: 비즈니스 모델
slide = add_content_slide(prs, "② 비즈니스 모델 (온/오프라인, 글로벌)", {})

models = [
    ("1. 오프라인 (B2G/B2B)", "관공서, 보건소, 대기업 ESG팀\n• 중증장애인생산품 우선구매\n• 기업 굿즈 납품"),
    ("2. 온라인 (D2C)", "가치소비 MZ세대\n• 자사몰, 스마트스토어\n• 와디즈/텀블벅 펀딩"),
    ("3. 글로벌 (Global)", "베트남 등 동남아 시장\n• 5년 차 해외 판로 개척\n• K-Eco 디자인 수출")
]

for i, (title, content) in enumerate(models):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(2.8), Inches(3))
    frame = box.text_frame
    frame.text = f"{title}\n\n{content}"
    frame.paragraphs[0].font.size = Pt(16)
    frame.paragraphs[0].font.bold = True
    for para in frame.paragraphs[1:]:
        para.font.size = Pt(13)

# Slide 5: 핵심 제품 소개
slide = add_content_slide(prs, "핵심 제품 소개 (Signature Products)", {})

box1 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(4), Inches(2.5))
frame1 = box1.text_frame
frame1.text = """1. 마블링 화분
폐플라스틱 고유의 패턴을 살린 세상에 하나뿐인 디자인.
기업 로고 각인 서비스 제공.
예상가: 12,000원

2. 아웃도어 카라비너
고강도 플라스틱(HDPE) 활용 캠핑/등산용 굿즈.
내구성 테스트 완료.
예상가: 3,500원"""
for para in frame1.paragraphs:
    para.font.size = Pt(14)

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3.5), Inches(2.5))
frame2 = box2.text_frame
frame2.text = """제품 차별성

자체 내구성 테스트 완료, 고품질 원료 사용으로 저가 제품 대비 우수한 품질.

제품 하단 'Made by 000(근로자 실명)' 각인 서비스."""
frame2.paragraphs[0].font.size = Pt(18)
frame2.paragraphs[0].font.bold = True
for para in frame2.paragraphs[1:]:
    para.font.size = Pt(14)

# Slide 6: SWOT 분석
slide = add_content_slide(prs, "④ 시장분석: SWOT 분석", {})

swot = [
    ("Strength (강점)", "• 복지 전문가+기술 전문가 시너지\n• 정부 인건비 지원 (가격 경쟁력)\n• 독창적인 금형 디자인 보유", BLUE),
    ("Weakness (약점)", "• 초기 브랜드 인지도 부족\n• 전용 설비 초기 투자비용 부담", RED),
    ("Opportunity (기회)", "• 공공기관 우선구매 시장(20조원)\n• 기업 ESG 경영 강화\n• 가치소비 트렌드 확산", GREEN),
    ("Threat (위협)", "• 저가 중국산 제품 경쟁\n• 유사 업사이클링 업체 난립", ORANGE)
]

positions = [(1, 2), (5.5, 2), (1, 4.5), (5.5, 4.5)]
for (title, content, color), (x, y) in zip(swot, positions):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.5), Inches(2))
    frame = box.text_frame
    frame.text = f"{title}\n\n{content}"
    frame.paragraphs[0].font.size = Pt(16)
    frame.paragraphs[0].font.bold = True
    frame.paragraphs[0].font.color.rgb = color
    for para in frame.paragraphs[1:]:
        para.font.size = Pt(13)

# Slide 7: 마케팅 전략
slide = add_content_slide(prs, "④ 마케팅 전략 및 판로 제시", {})

box1 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(4), Inches(3))
frame1 = box1.text_frame
frame1.text = """홍보 전략
"제품이 아닌 가치를 팝니다"

• 스토리텔링: 제품 QR코드로 제작자의 작업 영상 연결
• 체험단 운영: 지역 맘카페 연계 '나만의 화분 만들기'"""
frame1.paragraphs[0].font.size = Pt(18)
frame1.paragraphs[0].font.bold = True
for para in frame1.paragraphs[1:]:
    para.font.size = Pt(14)

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(4))
frame2 = box2.text_frame
frame2.text = """단계별 판로 개척

1단계 (공공): 나라장터, 꿈드래 쇼핑몰 입점

2단계 (온라인): 네이버 스마트스토어, 와디즈 펀딩

3단계 (제휴): 제로웨이스트 샵 20개소 입점

4단계 (수출): KOTRA 연계 해외 전시회 참가"""
frame2.paragraphs[0].font.size = Pt(18)
frame2.paragraphs[0].font.bold = True
for para in frame2.paragraphs[1:]:
    para.font.size = Pt(13)

# Slide 8: 인력 구성
slide = add_content_slide(prs, "③ 사업계획: 인력구성 (HR)", {})

roles = [
    ("대표 민소은 (본인)", "사회복지 전문가 (1급)", "인력 관리, 지자체/관공서 영업 총괄"),
    ("기술 이사 (CTO)", "생산 총괄 (경력 20년)", "사출 금형 설계, 생산 라인 및 품질 관리"),
    ("현장 근로자 (3명)", "취약계층 우선 채용", "경력단절여성 및 장애인\n사회복지관 추천 채용")
]

for i, (title, role, desc) in enumerate(roles):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.8), Inches(3))
    frame = box.text_frame
    frame.text = f"{title}\n\n{role}\n\n{desc}"
    frame.paragraphs[0].font.size = Pt(15)
    frame.paragraphs[0].font.bold = True
    for para in frame.paragraphs[1:]:
        para.font.size = Pt(12)

# Slide 9: 임금 계획
slide = add_content_slide(prs, "③ 사업계획: 인력 운영 및 임금 (2025년 기준)", {})

table_data = [
    ["구분", "대상", "월 급여 (세전)", "비고"],
    ["대표", "대표 민소은", "무급", "초기 3년간 무급 (재투자)"],
    ["관리직", "기술이사", "3,000,000원", "경력직 대우"],
    ["현장직", "취약계층 근로자", "2,096,270원", "최저임금 준수"]
]

from pptx.oxml.xmlchemy import OxmlElement

table = slide.shapes.add_table(len(table_data), 4, Inches(1), Inches(2), Inches(8), Inches(2)).table

for i, row_data in enumerate(table_data):
    for j, cell_value in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_value
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            if i == 0:
                paragraph.font.bold = True

note_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
note_frame = note_box.text_frame
note_frame.text = """※ 임금 산출 근거: 2025년 최저시급 10,030원 × 209시간
※ 복리후생: 4대보험 가입(두루누리), 식대 별도, 심리상담(EAP) 지원"""
for para in note_frame.paragraphs:
    para.font.size = Pt(12)

# Slide 10: 정부 제도 활용
slide = add_content_slide(prs, "③ 사업계획: 정부 제도 활용 전략", {})

stages = [
    ("1. 진입기", "사회적기업가 육성사업", "창업 초기 자금 3,000만원 확보하여 핵심 설비(사출기) 구입"),
    ("2. 성장기", "일자리창출사업 (예비)", "예비사회적기업 지정 후, 취약계층 신규 채용 인건비의 50~70% 지원"),
    ("3. 도약기", "장애인 표준사업장 지원", "한국장애인고용공단 시설 무상지원금으로 작업장 환경 개선")
]

for i, (stage, title, desc) in enumerate(stages):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.8), Inches(3))
    frame = box.text_frame
    frame.text = f"{stage}\n\n{title}\n\n{desc}"
    frame.paragraphs[0].font.size = Pt(14)
    frame.paragraphs[0].font.bold = True
    frame.paragraphs[0].font.color.rgb = BLUE
    for para in frame.paragraphs[1:]:
        para.font.size = Pt(12)

# Slide 11: 초기 투자 계획
slide = add_content_slide(prs, "초기 투자 및 생산 설비 계획", {})

box1 = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(3.5), Inches(3))
frame1 = box1.text_frame
frame1.text = """총 소요 예산: 1억 원

• 시설자금 (6,000만원):
  사출기, 분쇄기, 금형 제작비, 공장 보증금

• 운전자금 (4,000만원):
  초기 6개월 인건비, 재료비, 시제품 홍보비"""
frame1.paragraphs[0].font.size = Pt(18)
frame1.paragraphs[0].font.bold = True
for para in frame1.paragraphs[1:]:
    para.font.size = Pt(14)

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3.5), Inches(3))
frame2 = box2.text_frame
frame2.text = """생산 설비 특징

• 소형 사출기:
  다품종 소량 생산 용이, 장애인 접근성 고려

• 저소음 분쇄기:
  작업자 청력 보호를 위한 방음 박스 및 안전 센서"""
frame2.paragraphs[0].font.size = Pt(18)
frame2.paragraphs[0].font.bold = True
for para in frame2.paragraphs[1:]:
    para.font.size = Pt(14)

# Slide 12: 자금 조달
slide = add_content_slide(prs, "자금 조달 방안", {})

subtitle = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(0.5))
subtitle.text_frame.text = "전략: 상환 부담 없는 정부지원금 비율(50%) 확대"
subtitle.text_frame.paragraphs[0].font.size = Pt(16)
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

funds = [
    ("자기자본 (20%)", "2,000만 원", "대표자 출자 (책임경영)"),
    ("정부지원 (50%)", "5,000만 원", "사회적기업 육성사업\n장애인공단 기기지원"),
    ("정책융자 (30%)", "3,000만 원", "서민금융진흥원 (저리)")
]

for i, (title, amount, desc) in enumerate(funds):
    x = 0.8 + (i * 3)
    box = slide.shapes.add_textbox(Inches(x), Inches(2.8), Inches(2.8), Inches(2.5))
    frame = box.text_frame
    frame.text = f"{title}\n\n{amount}\n\n{desc}"
    frame.paragraphs[0].font.size = Pt(15)
    frame.paragraphs[0].font.bold = True
    frame.paragraphs[1].font.size = Pt(20)
    frame.paragraphs[1].font.bold = True
    frame.paragraphs[1].font.color.rgb = GREEN
    for para in frame.paragraphs[2:]:
        para.font.size = Pt(12)

# Slide 13: 5년 계획
slide = add_content_slide(prs, "⑤ 향후 추진 계획 (5년 계획)", {})

years = [
    ("1년차 (진입)", "• 법인 설립\n• 예비사회적기업 지정\n• 시제품 5종 개발"),
    ("2년차 (성장)", "• 매출 3억 달성\n• 장애인 표준사업장 신청 준비\n• 고정 거래처 10곳"),
    ("3년차 (도약)", "• 사회적기업 본인증\n• 장애인 표준사업장 인증 완료\n• 공장 확장 이전\n• 취약계층 10명 고용"),
    ("4년차 (확장)", "• 신제품 라인업 확대\n• 지역 협력 공장 네트워크 구축\n• 전국 유통망 확대"),
    ("5년차 (Global)", "• 해외 수출 시작(베트남)\n• 연 매출 10억 달성\n• K-Eco 브랜드화")
]

for i, (title, content) in enumerate(years):
    x = 0.5 + (i * 1.9)
    box = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(1.7), Inches(4))
    frame = box.text_frame
    frame.text = f"{title}\n\n{content}"
    frame.paragraphs[0].font.size = Pt(13)
    frame.paragraphs[0].font.bold = True
    if i < 2:
        frame.paragraphs[0].font.color.rgb = BLUE
    elif i < 4:
        frame.paragraphs[0].font.color.rgb = GREEN
    else:
        frame.paragraphs[0].font.color.rgb = ORANGE
    for para in frame.paragraphs[1:]:
        para.font.size = Pt(11)

# Slide 14: 추정 손익 계산서
slide = add_content_slide(prs, "추정 손익 계산서 (3개년)", {})

financial_data = [
    ["구분 (단위: 백만 원)", "1년차", "2년차", "3년차"],
    ["매출액", "120", "300", "500"],
    ["매출원가", "40", "100", "165"],
    ["판관비 (인건비 포함)", "100", "180", "250"],
    ["영업이익", "△20", "20", "85"]
]

table = slide.shapes.add_table(len(financial_data), 4, Inches(1.5), Inches(2), Inches(7), Inches(2.5)).table

for i, row_data in enumerate(financial_data):
    for j, cell_value in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_value
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            if i == 0 or i == 4:
                paragraph.font.bold = True

note_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(0.5))
note_frame = note_box.text_frame
note_frame.text = "* 1차년도 적자는 정부지원금(영업외수익)으로 보전하여 현금 흐름 유지."
note_frame.paragraphs[0].font.size = Pt(11)

# Slide 15: 손익분기점
slide = add_content_slide(prs, "⑤ 향후 추진 계획: 손익분기점 (BEP)", {})

box1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(3.5), Inches(2))
frame1 = box1.text_frame
frame1.text = "BEP 달성 예상 시점\n\n8개월"
frame1.paragraphs[0].font.size = Pt(16)
frame1.paragraphs[1].font.size = Pt(60)
frame1.paragraphs[1].font.bold = True
frame1.paragraphs[1].font.color.rgb = GREEN
frame1.paragraphs[1].alignment = PP_ALIGN.CENTER

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(4))
frame2 = box2.text_frame
frame2.text = """달성 조건 및 전략

• 달성 조건: 월 매출 1,400만 원
  (BEP = 고정비 700만원 ÷ 공헌이익률 50%)

• 고정비 최소화: 월 700만 원
  (인건비 지원 반영 후 자부담분)

• 공헌이익률: 50% (제조업 특성)

• 전략: B2B 대형 계약(월 500만 원 x 2건)
  수주 시 조기 달성 가능"""
frame2.paragraphs[0].font.size = Pt(16)
frame2.paragraphs[0].font.bold = True
for para in frame2.paragraphs[1:]:
    para.font.size = Pt(13)

# Slide 16: 사회적 가치
slide = add_content_slide(prs, "사회적 가치 측정 (SROI)", {})

box1 = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(3.5), Inches(3.5))
frame1 = box1.text_frame
frame1.text = """정량적 성과 (Quantitative)

• 연간 폐플라스틱 25톤 재활용
  (3년차 기준)

• 탄소 감축 효과
  (소나무 1,200그루 식재 상당)

• 취약계층 10명 일자리 창출"""
frame1.paragraphs[0].font.size = Pt(16)
frame1.paragraphs[0].font.bold = True
frame1.paragraphs[0].font.color.rgb = GREEN
for para in frame1.paragraphs[1:]:
    para.font.size = Pt(13)

box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3.5), Inches(3.5))
frame2 = box2.text_frame
frame2.text = """정성적 성과 (Qualitative)

• 근로자 우울감 척도 30% 개선

• 경제적 자립을 통한 탈수급
  (수급자 탈피)

• 지역사회 '환경 인식' 개선"""
frame2.paragraphs[0].font.size = Pt(16)
frame2.paragraphs[0].font.bold = True
frame2.paragraphs[0].font.color.rgb = BLUE
for para in frame2.paragraphs[1:]:
    para.font.size = Pt(13)

# Slide 17: Closing
add_title_slide(
    prs,
    "Closing",
    "",
    ""
)

# Add closing message
closing_slide = prs.slides[-1]
msg_box = closing_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
msg_frame = msg_box.text_frame
msg_frame.text = '''"Re:Spring은 단순한 공장이 아닙니다.
우리 지역사회의 가장 아픈 곳을 치유하는
'생산적 복지'의 거점입니다."

준비된 대표, 검증된 기술, 확실한 수요처.
이제 지원금이라는 마중물이 필요합니다.'''
for para in msg_frame.paragraphs:
    para.font.size = Pt(20)
    para.font.color.rgb = RGBColor(255, 255, 255)
    para.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('/home/user/claude-quickstarts/ReSpring_Business_Plan.pptx')
print("✓ PPT 파일 생성 완료: ReSpring_Business_Plan.pptx")
